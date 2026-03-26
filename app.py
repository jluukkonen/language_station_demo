import streamlit as st
import pandas as pd
import fitz  # PyMuPDF
import html
import re
from io import BytesIO
from llm_engine import process_text
from database import MOCK_STUDENTS, MOCK_COURSES, SAMPLE_TEXTS

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Page Config
st.set_page_config(
    page_title="Language Station",
    page_icon=None,
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for premium hackathon demo look
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');

    :root {
        --text-primary: #152033;
        --text-secondary: #5c6b82;
        --text-muted: #6b7a90;
        --bg-main: #f4f7fb;
        --bg-panel: rgba(255, 255, 255, 0.82);
        --bg-card: #ffffff;
        --bg-soft-blue: #eff5ff;
        --bg-soft-green: #effaf4;
        --bg-soft-ink: #eef2f8;
        --border-soft: rgba(22, 35, 58, 0.09);
        --surface-elevated: rgba(255, 255, 255, 0.92);
        --surface-interactive: rgba(255, 255, 255, 0.76);
        --sidebar-bg: linear-gradient(180deg, #f8fbff 0%, #eef4fb 100%);
        --hero-gradient: linear-gradient(135deg, #173870 0%, #2154b7 42%, #0b8cad 100%);
        --hero-glow: rgba(255, 255, 255, 0.38);
        --hero-chip-bg: rgba(255, 255, 255, 0.13);
        --hero-chip-border: rgba(255, 255, 255, 0.18);
        --hero-copy-color: rgba(255, 255, 255, 0.88);
        --success-border: rgba(31, 157, 99, 0.2);
        --info-border: rgba(41, 94, 239, 0.2);
        --success-text: #1f9d63;
        --warning-text: #7a5600;
        --scrollbar-thumb: #cdd5e0;
        --shadow-soft: 0 12px 30px rgba(23, 37, 84, 0.07);
        --shadow-hover: 0 18px 42px rgba(23, 37, 84, 0.11);
        --blue: #295eef;
        --blue-deep: #183b8c;
        --teal: #0e8aa8;
        --green: #1f9d63;
        --gold: #d9a441;
        --radius-lg: 22px;
        --radius-md: 16px;
        --radius-sm: 12px;
    }

    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
        color: var(--text-primary);
    }

    [data-testid="stAppViewContainer"] {
        background: var(--bg-main);
        background-image:
            radial-gradient(circle at top left, rgba(41, 94, 239, 0.12), transparent 35%),
            radial-gradient(circle at top right, rgba(14, 138, 168, 0.10), transparent 30%);
    }

    .main .block-container {
        padding-top: 1.7rem;
        padding-bottom: 2rem;
        max-width: 1280px;
    }

    [data-testid="stSidebar"] {
        background: var(--sidebar-bg);
        border-right: 1px solid var(--border-soft);
    }

    [data-testid="stSidebar"] > div:first-child {
        padding-top: 1.4rem;
    }

    .sidebar-shell {
        padding: 0.1rem 0 0.45rem 0;
    }

    .sidebar-brand {
        background: linear-gradient(180deg, rgba(255, 255, 255, 0.9), rgba(244, 248, 255, 0.92));
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        padding: 0.9rem 0.95rem;
        box-shadow: var(--shadow-soft);
        margin-bottom: 0.8rem;
    }

    .sidebar-kicker {
        font-size: 0.72rem;
        font-weight: 700;
        letter-spacing: 0.14em;
        text-transform: uppercase;
        color: var(--blue-deep);
        margin-bottom: 0.5rem;
    }

    .sidebar-title {
        font-size: 1.15rem;
        line-height: 1.2;
        font-weight: 800;
        color: var(--text-primary);
        margin-bottom: 0.35rem;
    }

    .sidebar-copy {
        font-size: 0.9rem;
        line-height: 1.5;
        color: var(--text-secondary);
        margin: 0;
    }

    .sidebar-section {
        margin: 0.95rem 0 0.55rem 0;
        font-size: 0.72rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 0.14em;
        color: var(--text-muted);
    }

    .sidebar-mini-card {
        background: linear-gradient(180deg, #ffffff 0%, #f8fbff 100%);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        box-shadow: var(--shadow-soft);
        padding: 0.9rem 0.95rem;
        margin-top: 0.55rem;
    }

    .sidebar-mini-title {
        font-size: 0.68rem;
        font-weight: 800;
        color: var(--text-muted);
        text-transform: uppercase;
        letter-spacing: 0.12em;
        margin-bottom: 0.55rem;
    }

    .sidebar-avatar-row {
        display: flex;
        align-items: center;
        gap: 0.1rem;
        margin-bottom: 0.55rem;
    }

    .sidebar-avatar {
        font-size: 1.08rem;
        line-height: 1;
        margin-right: -0.18rem;
    }

    .sidebar-meta-strong {
        font-size: 0.9rem;
        font-weight: 800;
        color: var(--text-primary);
        margin-bottom: 0.35rem;
    }

    .sidebar-chip-row {
        display: flex;
        flex-wrap: wrap;
        gap: 0.35rem;
        margin-bottom: 0.55rem;
    }

    .sidebar-chip {
        display: inline-flex;
        align-items: center;
        border-radius: 999px;
        padding: 0.18rem 0.48rem;
        font-size: 0.65rem;
        font-weight: 700;
        letter-spacing: 0.06em;
        text-transform: uppercase;
        border: 1px solid var(--border-soft);
    }

    .sidebar-chip.blue {
        background: var(--bg-soft-blue);
        color: var(--blue-deep);
    }

    .sidebar-chip.ink {
        background: var(--bg-soft-ink);
        color: var(--text-secondary);
    }

    .sidebar-note {
        font-size: 0.76rem;
        line-height: 1.45;
        color: var(--text-secondary);
    }

    .sidebar-divider {
        height: 1px;
        background: var(--border-soft);
        margin: 1rem 0 0.85rem 0;
    }

    .student-card {
        background: linear-gradient(180deg, #ffffff 0%, #f9fbff 100%);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        box-shadow: var(--shadow-soft);
        padding: 0.95rem;
        margin-bottom: 0.8rem;
    }

    .student-header {
        display: flex;
        align-items: center;
        justify-content: space-between;
        gap: 0.75rem;
        margin-bottom: 0.7rem;
    }

    .student-id {
        display: flex;
        align-items: center;
        gap: 0.55rem;
        min-width: 0;
    }

    .student-avatar {
        width: 2rem;
        height: 2rem;
        border-radius: 999px;
        display: inline-flex;
        align-items: center;
        justify-content: center;
        background: var(--bg-soft-blue);
        border: 1px solid var(--border-soft);
        font-size: 1.05rem;
        flex-shrink: 0;
    }

    .student-name {
        margin: 0;
        color: var(--text-primary);
        font-size: 0.95rem;
        font-weight: 800;
        line-height: 1.2;
    }

    .student-program {
        margin: 0.14rem 0 0 0;
        color: var(--text-secondary);
        font-size: 0.78rem;
        line-height: 1.35;
    }

    .student-badge {
        display: inline-flex;
        align-items: center;
        border-radius: 999px;
        padding: 0.26rem 0.56rem;
        background: var(--bg-soft-green);
        color: var(--success-text);
        border: 1px solid var(--border-soft);
        font-size: 0.68rem;
        font-weight: 800;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        flex-shrink: 0;
    }

    .student-chip-wrap {
        display: flex;
        flex-wrap: wrap;
        gap: 0.35rem;
        margin-bottom: 0.7rem;
    }

    .student-chip {
        display: inline-flex;
        align-items: center;
        border-radius: 999px;
        padding: 0.18rem 0.5rem;
        background: var(--bg-soft-ink);
        color: var(--text-secondary);
        border: 1px solid var(--border-soft);
        font-size: 0.66rem;
        font-weight: 700;
        letter-spacing: 0.04em;
    }

    .student-grid {
        display: grid;
        grid-template-columns: repeat(2, minmax(0, 1fr));
        gap: 0.65rem;
        margin-bottom: 0.7rem;
    }

    .student-block {
        background: #fbfcff;
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-sm);
        padding: 0.7rem 0.75rem;
    }

    .student-block-title {
        margin: 0 0 0.28rem 0;
        color: var(--text-muted);
        font-size: 0.66rem;
        font-weight: 800;
        letter-spacing: 0.1em;
        text-transform: uppercase;
    }

    .student-block-copy {
        margin: 0;
        color: var(--text-primary);
        font-size: 0.82rem;
        line-height: 1.45;
    }

    .student-callout-row {
        display: grid;
        grid-template-columns: 1fr;
        gap: 0.55rem;
    }

    .student-callout {
        border-radius: var(--radius-sm);
        padding: 0.78rem 0.82rem;
        border: 1px solid var(--border-soft);
    }

    .student-callout.barrier {
        background: #fff8f2;
        border-color: rgba(217, 164, 65, 0.22);
    }

    .student-callout.goal {
        background: #f5fbf8;
        border-color: rgba(31, 157, 99, 0.18);
    }

    .hero-card {
        position: relative;
        overflow: hidden;
        padding: 1.55rem 1.65rem 1.35rem 1.65rem;
        border-radius: 24px;
        background:
            radial-gradient(circle at 90% 10%, var(--hero-glow), transparent 18%),
            var(--hero-gradient);
        box-shadow: 0 22px 52px rgba(24, 59, 140, 0.18);
        border: 1px solid rgba(255, 255, 255, 0.14);
        margin-bottom: 1.2rem;
        color: white;
    }

    .hero-card::after {
        content: "";
        position: absolute;
        inset: auto -80px -120px auto;
        width: 280px;
        height: 280px;
        border-radius: 50%;
        background: rgba(255, 255, 255, 0.08);
        filter: blur(6px);
    }

    .hero-kicker {
        display: inline-block;
        padding: 0.42rem 0.68rem;
        border-radius: 999px;
        background: var(--hero-chip-bg);
        border: 1px solid var(--hero-chip-border);
        font-size: 0.7rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 0.12em;
        margin-bottom: 0.85rem;
    }

    .hero-title {
        font-size: clamp(1.9rem, 3.2vw, 3rem);
        line-height: 1.08;
        font-weight: 800;
        margin: 0;
        max-width: 700px;
    }

    .hero-copy {
        margin: 0.8rem 0 1rem 0;
        font-size: 0.96rem;
        line-height: 1.65;
        max-width: 720px;
        color: var(--hero-copy-color);
    }

    .hero-badges {
        display: flex;
        flex-wrap: wrap;
        gap: 0.55rem;
    }

    .hero-badge {
        padding: 0.48rem 0.78rem;
        border-radius: 999px;
        background: var(--hero-chip-bg);
        border: 1px solid var(--hero-chip-border);
        font-size: 0.8rem;
        font-weight: 600;
        color: white;
    }

    .section-intro {
        background: var(--bg-panel);
        backdrop-filter: blur(14px);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-lg);
        box-shadow: var(--shadow-soft);
        padding: 1.05rem 1.15rem;
        margin-bottom: 0.85rem;
    }

    .result-context {
        display: flex;
        flex-wrap: wrap;
        gap: 0.55rem;
        margin: 0.15rem 0 1rem 0;
    }

    .result-context-chip {
        display: inline-flex;
        align-items: center;
        padding: 0.45rem 0.72rem;
        border-radius: 999px;
        background: var(--bg-panel);
        border: 1px solid var(--border-soft);
        color: var(--text-secondary);
        font-size: 0.78rem;
        font-weight: 600;
        box-shadow: var(--shadow-soft);
    }

    .section-kicker {
        font-size: 0.72rem;
        letter-spacing: 0.14em;
        text-transform: uppercase;
        font-weight: 700;
        color: var(--blue-deep);
        margin-bottom: 0.45rem;
    }

    .section-title {
        margin: 0;
        color: var(--text-primary);
        font-size: 1.35rem;
        font-weight: 800;
    }

    .section-copy {
        margin: 0.45rem 0 0 0;
        color: var(--text-secondary);
        line-height: 1.65;
        font-size: 0.98rem;
    }

    .stats-grid {
        display: grid;
        grid-template-columns: repeat(3, minmax(0, 1fr));
        gap: 1rem;
        margin: 1rem 0 1.4rem 0;
    }

    .stat-card {
        background: var(--bg-panel);
        backdrop-filter: blur(14px);
        border: 1px solid var(--border-soft);
        border-radius: 22px;
        box-shadow: var(--shadow-soft);
        padding: 1.15rem 1.2rem;
    }

    .stat-label {
        color: var(--text-secondary);
        font-size: clamp(0.66rem, 0.61rem + 0.22vw, 0.78rem);
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 0.12em;
        margin-bottom: 0.6rem;
        line-height: 1.35;
    }

    .stat-value {
        color: var(--text-primary);
        font-size: clamp(1.4rem, 1.08rem + 0.9vw, 2rem);
        line-height: 1;
        font-weight: 800;
    }

    .stat-subtle {
        margin-top: 0.45rem;
        color: var(--text-muted);
        font-size: 0.92rem;
    }

    .stTabs [data-baseweb="tab-list"] {
        gap: 0.85rem;
        padding: 0.25rem 0 0.6rem 0;
    }

    .stTabs [data-baseweb="tab"] {
        height: 50px;
        background: var(--surface-interactive);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        padding: 0 1rem;
        color: var(--text-secondary);
        font-weight: 600;
        transition: all 0.2s ease;
        box-shadow: none;
    }

    .stTabs [data-baseweb="tab"]:hover {
        color: var(--text-primary);
        border-color: rgba(41, 94, 239, 0.18);
        transform: translateY(-1px);
    }

    .stTabs [aria-selected="true"] {
        background: linear-gradient(180deg, var(--bg-card) 0%, var(--bg-soft-blue) 100%);
        color: var(--blue-deep);
        border-color: rgba(41, 94, 239, 0.22);
        box-shadow: var(--shadow-soft);
    }

    .glossary-card,
    .planner-card,
    .text-panel {
        background: var(--surface-elevated);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        box-shadow: var(--shadow-soft);
        transition: transform 0.22s ease, box-shadow 0.22s ease, border-color 0.22s ease;
    }

    .glossary-card:hover,
    .planner-card:hover,
    .text-panel:hover {
        transform: translateY(-3px);
        box-shadow: var(--shadow-hover);
        border-color: rgba(41, 94, 239, 0.14);
    }

    .glossary-card {
        padding: 1.15rem;
        margin-bottom: 0.9rem;
        background: var(--bg-card);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
    }

    .card-meta {
        display: flex;
        align-items: center;
        flex-wrap: wrap;
        gap: 0.55rem;
        margin-bottom: 0.8rem;
    }

    .term-chip,
    .translation-chip,
    .difficulty-chip,
    .planner-chip {
        display: inline-flex;
        align-items: center;
        border-radius: 999px;
        font-size: 0.8rem;
        font-weight: 700;
        padding: 0.45rem 0.72rem;
    }

    .term-chip {
        background: var(--bg-soft-blue);
        color: var(--blue-deep);
        border: 1px solid var(--border-soft);
    }

    .translation-chip {
        background: linear-gradient(180deg, #f7fcf9 0%, var(--bg-soft-green) 100%);
        color: var(--success-text);
        border: 1px solid var(--border-soft);
        font-weight: 600;
    }

    .difficulty-chip {
        background: var(--bg-soft-ink);
        color: var(--text-secondary);
        border: 1px solid var(--border-soft);
    }

    .planner-chip {
        background: var(--bg-soft-blue);
        color: var(--blue-deep);
        margin-bottom: 0.8rem;
    }

    .term-title {
        color: var(--text-primary);
        font-size: 1.3rem;
        line-height: 1.18;
        font-weight: 800;
        margin: 0 0 0.7rem 0;
    }

    .card-label {
        display: block;
        color: var(--blue-deep);
        font-size: 0.78rem;
        font-weight: 700;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        margin-bottom: 0.32rem;
    }

    .card-copy {
        color: var(--text-primary);
        line-height: 1.66;
        margin: 0 0 0.78rem 0;
        font-size: 0.95rem;
    }

    .cognitive-note {
        margin-top: 1rem;
        padding: 1rem 1.25rem;
        border-radius: 12px;
        background: var(--bg-soft-ink);
        border-left: 4px solid var(--gold);
        color: var(--text-primary);
        font-size: 0.9rem;
        position: relative;
        line-height: 1.65;
        max-width: 92%;
    }

    .cognitive-note::before {
        content: "Teacher Insight";
        display: block;
        font-weight: 800;
        text-transform: uppercase;
        font-size: 0.7rem;
        margin-bottom: 0.25rem;
        letter-spacing: 0.05em;
        color: var(--warning-text);
    }

    .bridge-grid {
        display: grid;
        grid-template-columns: repeat(2, minmax(0, 1fr));
        gap: 1rem;
    }

    .bridge-stats {
        display: grid;
        grid-template-columns: repeat(3, minmax(0, 1fr));
        gap: 0.8rem;
        margin: 0.35rem 0 1rem 0;
    }

    .bridge-stat {
        background: var(--surface-elevated);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        box-shadow: var(--shadow-soft);
        padding: 0.85rem 0.95rem;
    }

    .bridge-stat-label {
        margin: 0;
        color: var(--text-muted);
        font-size: 0.7rem;
        font-weight: 800;
        letter-spacing: 0.12em;
        text-transform: uppercase;
    }

    .bridge-stat-value {
        margin: 0.28rem 0 0 0;
        color: var(--text-primary);
        font-size: 1.18rem;
        font-weight: 800;
        line-height: 1.1;
    }

    .text-panel {
        height: 500px;
        overflow-y: auto;
        padding: 1.5rem;
    }

    .text-panel.original {
        background: var(--bg-card);
        border: 1px solid var(--border-soft);
    }

    .text-panel.simplified {
        border: 2px solid var(--success-border);
        background: var(--bg-card);
        box-shadow: 0 0 20px rgba(31, 157, 99, 0.05);
    }

    .text-panel::-webkit-scrollbar {
        width: 6px;
    }

    .text-panel::-webkit-scrollbar-thumb {
        background: var(--scrollbar-thumb);
        border-radius: 10px;
    }

    .panel-label {
        display: inline-block;
        border-radius: 999px;
        padding: 0.45rem 0.72rem;
        font-size: 0.78rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        margin-bottom: 0.85rem;
    }

    .panel-label.original {
        background: var(--bg-soft-ink);
        color: var(--text-secondary);
    }

    .panel-label.simplified {
        background: var(--bg-soft-green);
        color: var(--success-text);
    }

    .panel-title {
        margin: 0 0 0.6rem 0;
        color: var(--text-primary);
        font-size: 1.15rem;
        font-weight: 800;
    }

    .panel-copy {
        margin: 0;
        color: var(--text-secondary);
        line-height: 1.78;
        font-size: 0.98rem;
        white-space: pre-wrap;
    }

    .planner-card {
        padding: 1.1rem;
        margin-bottom: 0.9rem;
        background: var(--bg-card);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
    }

    .planner-title {
        margin: 0 0 0.45rem 0;
        color: var(--text-primary);
        font-size: 1.08rem;
        font-weight: 800;
    }

    .planner-copy {
        margin: 0;
        color: var(--text-secondary);
        line-height: 1.72;
        font-size: 0.98rem;
    }

    .activity-grid {
        display: grid;
        grid-template-columns: 1.1fr 1.7fr;
        gap: 1rem;
        margin-top: 0.2rem;
    }

    .activity-side {
        display: flex;
        flex-direction: column;
        gap: 0.65rem;
    }

    .activity-block {
        background: #f8fbff;
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-sm);
        padding: 0.82rem 0.9rem;
    }

    .activity-block.soft-green {
        background: #f5fcf8;
    }

    .activity-block.soft-gold {
        background: #fffaf0;
    }

    .activity-block-title {
        margin: 0 0 0.35rem 0;
        color: var(--blue-deep);
        font-size: 0.72rem;
        font-weight: 800;
        letter-spacing: 0.12em;
        text-transform: uppercase;
    }

    .activity-block-copy {
        margin: 0;
        color: var(--text-primary);
        font-size: 0.92rem;
        line-height: 1.55;
    }

    .activity-steps {
        background: #ffffff;
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-sm);
        padding: 0.9rem 0.95rem;
    }

    .activity-step-list {
        margin: 0;
        padding-left: 1.15rem;
        color: var(--text-secondary);
    }

    .activity-step-list li {
        margin-bottom: 0.6rem;
        padding-left: 0.15rem;
        line-height: 1.62;
    }

    .activity-step-list li:last-child {
        margin-bottom: 0;
    }

    .status-banner {
        background: var(--bg-soft-green);
        border: 1px solid var(--success-border);
        color: var(--success-text);
        border-radius: 18px;
        padding: 0.95rem 1rem;
        box-shadow: var(--shadow-soft);
        margin-bottom: 1rem;
        font-weight: 600;
    }

    .kpi-grid {
        display: grid;
        grid-template-columns: repeat(4, minmax(0, 1fr));
        gap: 1rem;
        margin: 0.9rem 0 1.1rem 0;
    }

    .kpi-tile {
        background: var(--bg-card);
        border-radius: var(--radius-md);
        border: 1px solid var(--border-soft);
        box-shadow: var(--shadow-soft);
        padding: 0.95rem 1rem;
    }

    .kpi-tile.accent-blue {
        background: var(--bg-soft-blue);
        border-color: var(--info-border);
    }

    .kpi-tile.accent-green {
        background: var(--bg-soft-green);
        border-color: var(--success-border);
    }

    .kpi-tile.accent-gold {
        background: rgba(217, 164, 65, 0.1);
        border-color: rgba(217, 164, 65, 0.3);
    }

    .kpi-label {
        margin: 0;
        color: var(--text-muted);
        font-size: 0.68rem;
        font-weight: 800;
        letter-spacing: 0.14em;
        text-transform: uppercase;
    }

    .kpi-value {
        margin: 0.28rem 0 0 0;
        color: var(--text-primary);
        font-size: clamp(1.55rem, 1.15rem + 0.8vw, 2rem);
        font-weight: 800;
        line-height: 1.05;
    }

    .kpi-copy {
        margin: 0.35rem 0 0 0;
        color: var(--text-secondary);
        font-size: 0.82rem;
        line-height: 1.4;
    }

    .empty-state {
        background: var(--bg-panel);
        backdrop-filter: blur(8px);
        border: 1px solid var(--border-soft);
        border-radius: 24px;
        box-shadow: var(--shadow-soft);
        padding: 1.4rem 1.5rem;
        color: var(--text-secondary);
        line-height: 1.7;
    }

    .footer-note {
        text-align: center;
        color: var(--text-muted);
        font-size: 0.9rem;
        padding: 0.5rem 0 0.2rem 0;
    }

    .floating-toggle-label {
        font-size: 0.68rem;
        font-weight: 800;
        letter-spacing: 0.12em;
        text-transform: uppercase;
        color: var(--text-muted);
        margin: 0.2rem 0 0.35rem 0;
        text-align: center;
    }

    div[data-testid="stColumn"]:has(.floating-toggle-label) {
        position: sticky;
        top: 5.8rem;
        align-self: flex-start;
    }

    div[data-testid="stColumn"]:has(.floating-toggle-label) [data-testid="stRadio"] {
        background: rgba(255, 255, 255, 0.88);
        border: 1px solid var(--border-soft);
        border-radius: 999px;
        box-shadow: var(--shadow-soft);
        padding: 0.35rem 0.45rem 0.45rem 0.45rem;
        margin-top: 0;
        backdrop-filter: blur(10px);
    }

    div[data-testid="stColumn"]:has(.floating-toggle-label) [data-testid="stRadio"] label {
        justify-content: center;
    }

    div[data-testid="stColumn"]:has(.floating-toggle-label) [data-testid="stRadio"] [role="radiogroup"] {
        gap: 0.15rem;
    }

    [data-testid="stFileUploader"],
    [data-testid="stTextArea"],
    [data-testid="stRadio"],
    [data-baseweb="select"],
    [data-testid="stMultiSelect"] {
        background: var(--surface-interactive);
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        padding: 0.4rem 0.65rem 0.55rem 0.65rem;
        box-shadow: inset 0 1px 0 rgba(255, 255, 255, 0.05);
        margin-bottom: 0.75rem;
    }

    [data-testid="stRadio"] {
        padding: 0.72rem 0.8rem 0.82rem 0.8rem;
    }

    [data-testid="stTextArea"] textarea {
        line-height: 1.65;
    }

    [data-testid="stSidebar"] [data-testid="stExpander"] {
        border: 1px solid var(--border-soft);
        border-radius: var(--radius-md);
        background: rgba(255, 255, 255, 0.6);
        overflow: hidden;
    }

    [data-testid="stSidebar"] [data-testid="stExpander"] details summary {
        padding: 0.55rem 0.75rem;
    }

    .stButton > button {
        width: 100%;
        border: none;
        border-radius: var(--radius-sm);
        padding: 0.82rem 1rem;
        font-weight: 800;
        color: white;
        background: var(--hero-gradient);
        box-shadow: 0 16px 28px rgba(24, 59, 140, 0.22);
        transition: transform 0.18s ease, box-shadow 0.18s ease, filter 0.18s ease;
    }

    /* Hide theme switcher from the main menu */
    [data-testid="stMainMenu"] ul > div:first-child,
    [data-testid="stMainMenu"] ul > li:first-child {
        display: none !important;
    }

    /* Hide 'Made with Streamlit' and other footer elements from menu */
    .stMainMenu div[class*="StyledMenuFooter"] {
        display: none !important;
    }

    .stButton > button:hover {
        transform: translateY(-1px);
        filter: brightness(1.02);
        box-shadow: 0 20px 32px rgba(24, 59, 140, 0.28);
    }

    [data-testid="stNotification"],
    [data-testid="stAlert"] {
        border-radius: 18px;
    }

    @media (max-width: 980px) {
        .stats-grid,
        .bridge-grid,
        .kpi-grid,
        .bridge-stats,
        .activity-grid {
            grid-template-columns: 1fr;
        }

        div[data-testid="stColumn"]:has(.floating-toggle-label) {
            position: static;
        }

        .hero-card {
            padding: 1.2rem;
        }
    }

    @media (max-width: 1200px) {
        .stat-subtle {
            font-size: 0.84rem;
        }

        .activity-grid {
            grid-template-columns: 1fr;
        }
    }
    </style>
    """, unsafe_allow_html=True)

def format_activity_sections(instructions: str, language_code: str = "en") -> dict:
    text = " ".join(instructions.split())
    fields = {
        "goal": "",
        "task": "",
        "steps": "",
        "languages": "",
        "materials": "",
        "why": "",
    }

    task_label = r"(?:Task(?: \([^)]+\))?|Tehtävä(?: \([^)]+\))?)"
    patterns = {
        "goal": rf"(?:^|\b)(?:Goal|Tavoite):\s*(.*?)(?=\b(?:{task_label}|Steps?|Vaiheet?|Languages used|Käytetyt kielet|Materials?|Materiaalit|Why this works|Miksi tämä toimii):|$)",
        "task": rf"(?:^|\b){task_label}:\s*(.*?)(?=\b(?:Goal|Tavoite|Steps?|Vaiheet?|Languages used|Käytetyt kielet|Materials?|Materiaalit|Why this works|Miksi tämä toimii):|$)",
        "steps": rf"(?:^|\b)(?:Steps?|Vaiheet?):\s*(.*?)(?=\b(?:Goal|Tavoite|{task_label}|Languages used|Käytetyt kielet|Materials?|Materiaalit|Why this works|Miksi tämä toimii):|$)",
        "languages": rf"(?:^|\b)(?:Languages used|Käytetyt kielet):\s*(.*?)(?=\b(?:Goal|Tavoite|{task_label}|Steps?|Vaiheet?|Materials?|Materiaalit|Why this works|Miksi tämä toimii):|$)",
        "materials": rf"(?:^|\b)(?:Materials?|Materiaalit):\s*(.*?)(?=\b(?:Goal|Tavoite|{task_label}|Steps?|Vaiheet?|Languages used|Käytetyt kielet|Why this works|Miksi tämä toimii):|$)",
        "why": rf"(?:^|\b)(?:Why this works|Miksi tämä toimii):\s*(.*?)(?=\b(?:Goal|Tavoite|{task_label}|Steps?|Vaiheet?|Languages used|Käytetyt kielet|Materials?|Materiaalit):|$)",
    }

    for key, pattern in patterns.items():
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            fields[key] = match.group(1).strip()

    if not fields["steps"]:
        fallback = fields["task"] or text
        step_candidates = re.split(r"\s*(?:(?<=\.)\s+|\b\d+\.\s+)", fallback)
        clean_steps = [s.strip(" .") for s in step_candidates if len(s.strip()) > 20]
        fields["step_list"] = clean_steps[:6] if clean_steps else [fallback]
    else:
        step_candidates = re.split(r"\s*(?:\b\d+\.\s+|(?<=\.)\s+)", fields["steps"])
        clean_steps = [s.strip(" .") for s in step_candidates if len(s.strip()) > 10]
        fields["step_list"] = clean_steps[:6] if clean_steps else [fields["steps"]]

    if language_code == "fi":
        if not fields["goal"]:
            fields["goal"] = "Käytä tehtävää tukemaan monikielistä merkityksenmuodostusta lähdetekstin keskeisten ajatusten ympärillä."
        if not fields["languages"]:
            fields["languages"] = "Käytä englantia lähdekäsitteisiin ja luokkakeskusteluun, mutta salli kotikielet selventämiseen ja yhteiseen merkityksenrakentamiseen."
        if not fields["why"]:
            fields["why"] = "Tämä tehtävä vähentää kielellistä kuormitusta jakamalla ymmärrystä vertaisille, eri kielille sekä suulliseen ja visuaaliseen työskentelyyn."
        if not fields["task"]:
            fields["task"] = "Ohjaa opiskelijoita tehtävän läpi käyttäen lähdetekstiä, sanastoa ja yhteistoiminnallista keskustelua."
    else:
        if not fields["goal"]:
            fields["goal"] = "Use the activity to support multilingual meaning-making around the key ideas in the source text."
        if not fields["languages"]:
            fields["languages"] = "Use English for source concepts and classroom discussion, while allowing home languages for clarification and collaborative meaning-making."
        if not fields["why"]:
            fields["why"] = "This activity lowers language load by distributing understanding across peers, languages, and visual or oral explanation."
        if not fields["task"]:
            fields["task"] = "Guide students through the activity using the source text, glossary, and collaborative discussion."

    return fields


def get_localized_text(value, language_code: str) -> str:
    """Support both new bilingual payloads and legacy single-language strings."""
    if hasattr(value, language_code):
        localized = getattr(value, language_code, "")
        fallback = getattr(value, "en", "") or getattr(value, "fi", "")
        return localized or fallback or ""
    if isinstance(value, dict):
        localized = value.get(language_code, "")
        fallback = value.get("en", "") or value.get("fi", "")
        return localized or fallback or ""
    return str(value or "")

# Header
st.markdown("""
    <section class="hero-card">
        <div class="hero-kicker">AI-Powered Multilingual Pedagogy</div>
        <h1 class="hero-title">Transform monolingual materials into multilingual learning experiences in seconds.</h1>
        <p class="hero-copy">
            Upload a paragraph, article, or PDF to generate a cognitive glossary, a CEFR-adapted bridge text,
            and classroom activities that reduce language load without losing academic rigor.
        </p>
        <div class="hero-badges">
            <span class="hero-badge">Cognitive Glossary</span>
            <span class="hero-badge">B1 Simplification</span>
            <span class="hero-badge">Translanguaging Activities</span>
            <span class="hero-badge">Gemini 2.5 Flash / Pro</span>
        </div>
    </section>
""", unsafe_allow_html=True)

# Sidebar
st.sidebar.markdown(f"""
    <div style="display: flex; align-items: center; gap: 10px; margin-bottom: 20px;">
        <div style="background: var(--blue-deep); width: 40px; height: 40px; border-radius: 8px; display: flex; align-items: center; justify-content: center; color: white; font-weight: 800;">LS</div>
        <div>
            <div style="font-size: 14px; font-weight: 800; color: var(--text-primary); line-height:1;">Language Station</div>
            <div style="font-size: 10px; color: var(--text-secondary);">UEF Multilingual Pedagogy</div>
        </div>
    </div>
""", unsafe_allow_html=True)
st.sidebar.markdown("""
    <div class="sidebar-shell">
        <div class="sidebar-brand">
            <div class="sidebar-kicker">Language Station</div>
            <div class="sidebar-title">Teacher Control Center</div>
            <p class="sidebar-copy">
                Prepare multilingual classroom support in a single flow with outputs ready for teaching use.
            </p>
        </div>
    </div>
""", unsafe_allow_html=True)
selected_model = "gemini-2.5-flash"

# Initialize session state for the result and source text
if 'result' not in st.session_state:
    st.session_state.result = None
if 'source_text' not in st.session_state:
    st.session_state.source_text = ""
if 'output_ui_language' not in st.session_state:
    st.session_state.output_ui_language = "EN"
if "source_processing_summary" not in st.session_state:
    st.session_state.source_processing_summary = None
if "batch_results" not in st.session_state:
    st.session_state.batch_results = None

try:
    with open("assets/demo_samples.txt", "r", encoding="utf-8") as f:
        demo_text = f.read()
except Exception:
    demo_text = ""

# The primary state for the input area
if "source_text_input" not in st.session_state:
    st.session_state.source_text_input = demo_text

st.sidebar.markdown("""<div class="sidebar-section">Teaching Setup</div>""", unsafe_allow_html=True)
input_mode = st.sidebar.radio("Input Type", ["Use course material", "Describe your lesson"])
source_material_mode = "Single file"
if input_mode == "Use course material":
    source_material_mode = st.sidebar.radio("Source Mode", ["Single file", "Course Pack"])

if input_mode == "Use course material" and source_material_mode == "Course Pack":
    language_direction = None
    st.sidebar.caption("Course Pack mode detects source language per file and applies the direction automatically.")
else:
    language_direction = st.sidebar.radio("Language Support", ["English → Finnish", "Finnish → English"])
selected_students = st.sidebar.multiselect(
    "Select classroom group",
    MOCK_STUDENTS,
    default=MOCK_STUDENTS[:2],
    format_func=lambda x: x["name"]
)

# Mandatory Course Context Selection
selected_course = st.sidebar.selectbox(
    "Select course context",
    MOCK_COURSES,
    format_func=lambda x: x["name"]
)

if not selected_students:
    st.sidebar.info("No students selected. Language Station will generate terminology and accessible text without learning activities.")

# Styled Group Profile Card
if selected_students:
    all_langs = set()
    for s in selected_students:
        all_langs.update(s.get('linguistic_repertoire', []))
    
    st.sidebar.markdown(f"""
        <div class="sidebar-mini-card">
            <div class="sidebar-mini-title">Classroom Group Summary</div>
            <div class="sidebar-avatar-row">
                {" ".join([f"<span class='sidebar-avatar'>{s['avatar']}</span>" for s in selected_students])}
            </div>
            <div class="sidebar-meta-strong">{len(selected_students)} student group</div>
            <div class="sidebar-chip-row">
                <span class="sidebar-chip blue">{len(all_langs)} languages</span>
                <span class="sidebar-chip ink">{", ".join(sorted(set(s['cefr'] for s in selected_students)))}</span>
            </div>
            <div class="sidebar-note"><strong>Linguistic repertoire:</strong> {", ".join(sorted(list(all_langs)))}</div>
        </div>
    """, unsafe_allow_html=True)

    with st.sidebar.expander("View student profiles", expanded=False):
        for student in selected_students:
            completed_courses = ", ".join(student.get("completed_courses", [])) or "None yet"
            linguistic_repertoire = ", ".join(student.get("linguistic_repertoire", [])) or "N/A"
            home_languages = ", ".join(student.get("home_languages", [])) or "N/A"
            student_name = html.escape(student["name"])
            student_avatar = html.escape(student["avatar"])
            student_program = html.escape(student["program"])
            student_cefr = html.escape(student["cefr"])
            student_year = html.escape(str(student["year"]))
            schooling_language = html.escape(student["schooling_language"])
            home_languages_safe = html.escape(home_languages)
            linguistic_repertoire_safe = html.escape(linguistic_repertoire)
            completed_courses_safe = html.escape(completed_courses)
            linguistic_barrier = html.escape(student["linguistic_barrier"])
            learning_goal = html.escape(student["goal"])

            st.markdown(f"""
                <div class="student-card">
                    <div class="student-header">
                        <div class="student-id">
                            <div class="student-avatar">{student_avatar}</div>
                            <div>
                                <p class="student-name">{student_name}</p>
                                <p class="student-program">Year {student_year} · {student_program}</p>
                            </div>
                        </div>
                        <div class="student-badge">{student_cefr}</div>
                    </div>
                    <div class="student-chip-wrap">
                        <span class="student-chip">Schooling: {schooling_language}</span>
                        <span class="student-chip">Home: {home_languages_safe}</span>
                    </div>
                    <div class="student-grid">
                        <div class="student-block">
                            <p class="student-block-title">Linguistic Repertoire</p>
                            <p class="student-block-copy">{linguistic_repertoire_safe}</p>
                        </div>
                        <div class="student-block">
                            <p class="student-block-title">Completed Courses</p>
                            <p class="student-block-copy">{completed_courses_safe}</p>
                        </div>
                    </div>
                    <div class="student-callout-row">
                        <div class="student-callout barrier">
                            <p class="student-block-title">Main Barrier</p>
                            <p class="student-block-copy">{linguistic_barrier}</p>
                        </div>
                        <div class="student-callout goal">
                            <p class="student-block-title">Learning Goal</p>
                            <p class="student-block-copy">{learning_goal}</p>
                        </div>
                    </div>
                </div>
            """, unsafe_allow_html=True)

st.sidebar.markdown("""
    <div class="sidebar-section" style="margin-top: 1.2rem;">Source Material</div>
""", unsafe_allow_html=True)
uploaded_file = None
uploaded_files = []
if input_mode == "Use course material":
    st.sidebar.markdown('<p style="font-size: 0.78rem; margin-bottom: 0.2rem; font-weight: 700; color: var(--text-secondary);">Sample materials</p>', unsafe_allow_html=True)
    
    # Minimalist language switch for samples
    sample_lang = st.sidebar.radio(
        "Sample language",
        ["ENG", "FIN"],
        horizontal=True,
        label_visibility="collapsed",
        format_func=lambda x: {"ENG": "🇬🇧", "FIN": "🇫🇮"}[x],
        key="sample_lang_choice"
    )
    st.sidebar.caption("Load a prepared example to preview the workflow quickly.")
    
    sample_keys = list(SAMPLE_TEXTS.keys())
    top_left, top_right = st.sidebar.columns(2)
    sample_button_layout = [
        ("Microbiology", 0, top_left),
        ("Public Health", 1, top_right),
    ]

    for btn_label, i, col in sample_button_layout:
        with col:
            if st.button(btn_label, key=f"btn_sample_{i}", use_container_width=True):
                lang_key = "en" if sample_lang == "ENG" else "fi"
                st.session_state.source_text_input = SAMPLE_TEXTS[sample_keys[i]][lang_key]
                st.rerun()

    if len(sample_keys) > 2:
        if st.sidebar.button("Biomedical Engineering", key="btn_sample_2", use_container_width=True):
            lang_key = "en" if sample_lang == "ENG" else "fi"
            st.session_state.source_text_input = SAMPLE_TEXTS[sample_keys[2]][lang_key]
            st.rerun()

    input_text = st.sidebar.text_area(
        "Paste your course material here:",
        height=200,
        placeholder="Paste a paragraph from a lecture, article, or course material...",
        key="source_text_input"
    )

    if source_material_mode == "Single file":
        uploaded_file = st.sidebar.file_uploader("Upload an academic text, PDF, or PowerPoint", type=["txt", "pdf", "pptx"])
    else:
        uploaded_files = st.sidebar.file_uploader(
            "Upload several course files",
            type=["txt", "pdf", "pptx"],
            accept_multiple_files=True,
            help="Batch process several course files independently. Each file gets its own detected language direction and result."
        )
else:
    # Uses the description from the top-level course selector
    input_text = st.sidebar.text_area(
        "Describe your lesson (topic, students, goals, materials):",
        height=300,
        value=selected_course["description"],
        placeholder="e.g., A lesson about handwashing for first-year nursing students. Focus on hygiene and microbiology terms."
    )

st.sidebar.markdown("### Generate Language Station")
generate_button = st.sidebar.button("Generate Language Station", type="primary")

# Data Attribution & Licensing (Small, muted text for compliance)
st.sidebar.markdown("""<div class="sidebar-divider"></div>""", unsafe_allow_html=True)
with st.sidebar.expander("Data attribution", expanded=False):
    st.markdown("""
        <div class="sidebar-mini-card" style="margin-top: 0;">
            <div class="sidebar-mini-title">Data Attribution Notice</div>
            <p style="font-size: 10px; color: var(--text-secondary); line-height: 1.5; margin-bottom: 10px;">
                This terminology dataset utilizes deterministic mappings to ensure clinical and academic reliability. It includes content from the following authoritative sources:
            </p>
            <ul style="font-size: 10px; color: var(--text-secondary); line-height: 1.5; padding-left: 14px; margin-bottom: 10px;">
                <li><strong>Finto.fi / National Library of Finland:</strong> YSO, KOKO, TERO, AFO, JUPO, KASSU, LAJISTO, and OIKO Ontologies. Licensed under <a href="https://creativecommons.org/licenses/by/4.0/" style="color: var(--text-secondary); text-decoration: underline;">CC BY 4.0</a>.</li>
                <li><strong>FinMeSH:</strong> Finnish translation of Medical Subject Headings. Produced by the National Library of Finland. Licensed under <a href="https://creativecommons.org/licenses/by/4.0/" style="color: var(--text-secondary); text-decoration: underline;">CC BY 4.0</a>. (Original MeSH data courtesy of the U.S. National Library of Medicine).</li>
                <li><strong>Tieteen termipankki / Sanastokeskus:</strong> Various domain-specific vocabularies. Licensed under <a href="https://creativecommons.org/licenses/by-sa/4.0/" style="color: var(--text-secondary); text-decoration: underline;">CC BY-SA 4.0</a>.</li>
            </ul>
            <p style="font-size: 10px; color: var(--text-secondary); font-style: italic; line-height: 1.4;">
                <strong>Disclaimer:</strong> This data is utilized as-is for terminology mapping purposes. Providers bear no responsibility for any modifications or applications of this tool.
            </p>
        </div>
    """, unsafe_allow_html=True)

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_text_from_pdf_bytes(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def _collect_shape_text(shape):
    collected_lines = []

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            runs_text = "".join(run.text for run in paragraph.runs).strip()
            paragraph_text = runs_text or paragraph.text.strip()
            if not paragraph_text:
                continue

            if getattr(paragraph, "level", 0) > 0:
                indent = "  " * paragraph.level
                collected_lines.append(f"{indent}- {paragraph_text}")
            else:
                collected_lines.append(paragraph_text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_values = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_values.append(cell_text)
            if row_values:
                collected_lines.append(" | ".join(row_values))

    return collected_lines


def _extract_speaker_notes(notes_slide):
    notes_lines = []
    if not notes_slide:
        return notes_lines

    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if not text:
            continue

        placeholder = getattr(shape, "placeholder_format", None)
        placeholder_type = str(getattr(placeholder, "type", "")).upper()
        if "SLIDE_IMAGE" in placeholder_type:
            continue

        notes_lines.append(text)

    return notes_lines


def extract_text_from_pptx(file_bytes):
    if Presentation is None:
        raise RuntimeError("PowerPoint support requires the python-pptx package.")

    presentation = Presentation(BytesIO(file_bytes))
    slide_blocks = []
    slides_with_notes = 0
    visual_warning = False

    for index, slide in enumerate(presentation.slides, start=1):
        slide_title = ""
        slide_lines = []

        title_shape = getattr(slide.shapes, "title", None)
        if title_shape and title_shape.text:
            slide_title = title_shape.text.strip()

        for shape in slide.shapes:
            slide_lines.extend(_collect_shape_text(shape))

        if slide_title:
            deduped_lines = []
            for line in slide_lines:
                if line.strip() == slide_title:
                    continue
                deduped_lines.append(line)
            slide_lines = deduped_lines

        notes_lines = _extract_speaker_notes(getattr(slide, "notes_slide", None))
        if notes_lines:
            slides_with_notes += 1

        if not slide_lines and not notes_lines:
            visual_warning = True
            slide_lines = ["This slide may rely on visual content not captured automatically."]

        slide_header = f"Slide {index}: {slide_title}" if slide_title else f"Slide {index}"
        block_parts = [slide_header]
        block_parts.extend(slide_lines)
        if notes_lines:
            block_parts.append("Speaker notes:")
            block_parts.extend(notes_lines)

        slide_blocks.append("\n".join(part for part in block_parts if part.strip()))

    extracted_text = "\n\n".join(slide_blocks).strip()
    return extracted_text, {
        "kind": "pptx",
        "slides_processed": len(presentation.slides),
        "slides_with_notes": slides_with_notes,
        "has_notes": slides_with_notes > 0,
        "visual_warning": visual_warning,
    }


def truncate_pptx_text(extracted_text, limit=15000):
    if len(extracted_text) <= limit:
        return extracted_text, False

    slide_blocks = [block.strip() for block in extracted_text.split("\n\n") if block.strip()]
    truncated_blocks = []
    used = 0

    for block in slide_blocks:
        lines = [line for line in block.splitlines() if line.strip()]
        if not lines:
            continue

        header = lines[0]
        visible_lines = []
        notes_lines = []
        in_notes = False
        for line in lines[1:]:
            if line.strip() == "Speaker notes:":
                in_notes = True
                continue
            if in_notes:
                notes_lines.append(line)
            else:
                visible_lines.append(line)

        candidate_parts = [header]
        candidate_parts.extend(visible_lines)
        if notes_lines:
            candidate_parts.append("Speaker notes:")
            candidate_parts.extend(notes_lines)

        for cutoff in (len(candidate_parts), 1 + len(visible_lines), 1):
            parts = candidate_parts[:cutoff]
            block_text = "\n".join(parts).strip()
            if not block_text:
                continue

            separator = "\n\n" if truncated_blocks else ""
            projected = used + len(separator) + len(block_text)
            if projected <= limit:
                truncated_blocks.append(block_text)
                used = projected
                break
        else:
            remaining = limit - used - (2 if truncated_blocks else 0)
            if remaining <= len(header):
                break
            shortened = block[:remaining].rstrip()
            if shortened:
                if truncated_blocks:
                    truncated_blocks.append(shortened)
                else:
                    truncated_blocks = [shortened]
            break

    return "\n\n".join(truncated_blocks).strip(), True


def detect_source_language(text: str):
    normalized = " ".join((text or "").split()).lower()
    tokens = re.findall(r"[a-zA-ZÀ-ÿäöåÄÖÅ]+", normalized)
    if len(tokens) < 25:
        return {
            "source_language": "Needs review",
            "confidence": "low",
            "resolved_direction": None,
            "reason": "Too little text to confidently detect language.",
        }

    english_markers = {
        "the", "and", "is", "are", "with", "for", "from", "this", "that", "students",
        "learning", "course", "teacher", "language", "text", "use", "can", "will", "into",
    }
    finnish_markers = {
        "ja", "on", "ovat", "että", "sekä", "tämä", "opiskelijat", "kieli", "oppiminen",
        "kurssi", "opettaja", "teksti", "voi", "ovatko", "suomen", "englannin", "kanssa",
    }

    english_score = sum(1 for token in tokens if token in english_markers)
    finnish_score = sum(1 for token in tokens if token in finnish_markers)
    finnish_score += normalized.count("ä") + normalized.count("ö") + normalized.count("å")

    total_score = english_score + finnish_score
    if total_score < 3:
        return {
            "source_language": "Needs review",
            "confidence": "low",
            "resolved_direction": None,
            "reason": "Language detection confidence is too low for this file.",
        }

    difference = abs(english_score - finnish_score)
    confidence = "high" if difference >= 3 else "low"
    if confidence == "low":
        return {
            "source_language": "Needs review",
            "confidence": confidence,
            "resolved_direction": None,
            "reason": "The file appears mixed-language or ambiguous.",
        }

    source_language = "English" if english_score > finnish_score else "Finnish"
    resolved_direction = "English → Finnish" if source_language == "English" else "Finnish → English"
    return {
        "source_language": source_language,
        "confidence": confidence,
        "resolved_direction": resolved_direction,
        "reason": f"Detected {source_language.lower()} source content.",
    }


def extract_source_material(file_obj):
    file_bytes = file_obj.getvalue()
    uploaded_name = file_obj.name.lower()
    processing_summary = {"kind": uploaded_name.rsplit(".", 1)[-1] if "." in uploaded_name else "unknown"}

    if uploaded_name.endswith(".pdf"):
        source_text = extract_text_from_pdf_bytes(file_bytes)
    elif uploaded_name.endswith(".pptx"):
        source_text, processing_summary = extract_text_from_pptx(file_bytes)
        source_text, was_truncated = truncate_pptx_text(source_text)
        processing_summary["was_truncated"] = was_truncated
    elif uploaded_name.endswith(".txt"):
        source_text = file_bytes.decode("utf-8")
    else:
        raise ValueError("Unsupported file type.")

    if "was_truncated" not in processing_summary:
        processing_summary["was_truncated"] = len(source_text) > 15000

    return source_text, processing_summary


def build_course_pack_preview(files):
    preview_rows = []
    for file_obj in files:
        filename = file_obj.name
        extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else "unknown"
        try:
            source_text, processing_summary = extract_source_material(file_obj)
            detection = detect_source_language(source_text)
            status = "Ready" if detection["resolved_direction"] else "Needs review"
            preview_rows.append({
                "File": filename,
                "Type": extension.upper(),
                "Detected": detection["source_language"],
                "Direction": detection["resolved_direction"] or "Skipped",
                "Status": status,
                "Notes": detection["reason"],
            })
        except Exception as exc:
            preview_rows.append({
                "File": filename,
                "Type": extension.upper(),
                "Detected": "Unreadable",
                "Direction": "Skipped",
                "Status": "Error",
                "Notes": str(exc),
            })
    return preview_rows

# Main Content Logic
if generate_button:
    if input_mode == "Use course material" and source_material_mode == "Course Pack":
        if not uploaded_files:
            st.warning("Please upload one or more course files to run Course Pack mode.")
        else:
            with st.spinner("Processing course files one by one and detecting the correct language direction for each file..."):
                batch_results = []
                for file_obj in uploaded_files:
                    extension = file_obj.name.rsplit(".", 1)[-1].lower() if "." in file_obj.name else "unknown"
                    try:
                        source_text, processing_summary = extract_source_material(file_obj)
                        detection = detect_source_language(source_text)

                        if not source_text.strip():
                            batch_results.append({
                                "filename": file_obj.name,
                                "file_type": extension,
                                "source_text": "",
                                "source_language": "Needs review",
                                "resolved_direction": None,
                                "status": "Skipped",
                                "warnings": ["No extractable text found in this file."],
                                "processing_summary": processing_summary,
                                "result": None,
                            })
                            continue

                        if not detection["resolved_direction"]:
                            batch_results.append({
                                "filename": file_obj.name,
                                "file_type": extension,
                                "source_text": source_text,
                                "source_language": detection["source_language"],
                                "resolved_direction": None,
                                "status": "Needs review",
                                "warnings": [detection["reason"]],
                                "processing_summary": processing_summary,
                                "result": None,
                            })
                            continue

                        result = process_text(
                            source_text,
                            model_type=selected_model,
                            input_mode="Use course material",
                            language_direction=detection["resolved_direction"],
                            selected_students=selected_students,
                            selected_course=selected_course
                        )

                        warnings = []
                        if processing_summary.get("was_truncated"):
                            warnings.append("This file was truncated to fit the analysis limit.")
                        if processing_summary.get("visual_warning"):
                            warnings.append("Some presentation slides may rely on visuals that were not fully captured.")

                        batch_results.append({
                            "filename": file_obj.name,
                            "file_type": extension,
                            "source_text": source_text,
                            "source_language": detection["source_language"],
                            "resolved_direction": detection["resolved_direction"],
                            "status": "Processed",
                            "warnings": warnings,
                            "processing_summary": processing_summary,
                            "result": result,
                        })
                    except Exception as e:
                        batch_results.append({
                            "filename": file_obj.name,
                            "file_type": extension,
                            "source_text": "",
                            "source_language": "Unreadable",
                            "resolved_direction": None,
                            "status": "Error",
                            "warnings": [str(e)],
                            "processing_summary": {"kind": extension},
                            "result": None,
                        })

                st.session_state.batch_results = batch_results
                st.session_state.result = None
                st.session_state.source_processing_summary = None
    else:
        source_text = ""
        processing_summary = None
        if uploaded_file:
            source_text, processing_summary = extract_source_material(uploaded_file)
        elif input_text:
            source_text = input_text
        
        if not source_text.strip():
            st.warning("Please enter or upload an academic text to generate outputs.")
        else:
            # UI Warning for truncation
            was_truncated = len(source_text) > 15000
            if processing_summary is None and was_truncated:
                st.info("Note: The source text is very long. We have truncated it to the most relevant first 15,000 characters for the analysis.")
            elif processing_summary and processing_summary.get("was_truncated"):
                st.info("Note: This presentation was condensed to fit the first 15,000 characters while preserving slide coverage.")

            with st.spinner("Analyzing linguistic complexity and generating multilingual learning resources..."):
                try:
                    # Call the LLM engine and store in session state
                    st.session_state.result = process_text(
                        source_text, 
                        model_type=selected_model,
                        input_mode=input_mode,
                        language_direction=language_direction,
                        selected_students=selected_students,
                        selected_course=selected_course
                    )
                    st.session_state.source_text_cache = source_text # Cache source text too
                    st.session_state.lang_dir_cache = language_direction # Cache language direction
                    st.session_state.source_processing_summary = processing_summary
                    st.session_state.batch_results = None
                except Exception as e:
                    st.error(f"An error occurred during processing: {str(e)}")

if input_mode == "Use course material" and source_material_mode == "Course Pack" and uploaded_files:
    preview_rows = build_course_pack_preview(uploaded_files)
    st.sidebar.markdown("#### Course Pack Preview")
    st.sidebar.dataframe(pd.DataFrame(preview_rows), use_container_width=True, hide_index=True)

# Display Logic (Outside the button click block, triggered by session state)
if st.session_state.batch_results and input_mode == "Use course material" and source_material_mode == "Course Pack":
    batch_results = st.session_state.batch_results
    has_selected_students = bool(selected_students)
    processed_count = sum(1 for item in batch_results if item["status"] == "Processed")
    skipped_count = sum(1 for item in batch_results if item["status"] in {"Skipped", "Needs review", "Error"})
    english_count = sum(1 for item in batch_results if item["source_language"] == "English")
    finnish_count = sum(1 for item in batch_results if item["source_language"] == "Finnish")

    results_col, toggle_col = st.columns([0.92, 0.08], gap="medium")
    with toggle_col:
        st.markdown('<div class="floating-toggle-label">Output</div>', unsafe_allow_html=True)
        output_language = st.radio(
            "AI output language",
            ["EN", "FIN"],
            horizontal=False,
            label_visibility="collapsed",
            format_func=lambda x: {"EN": "🇬🇧", "FIN": "🇫🇮"}[x],
            key="output_ui_language_batch"
        )

    output_language_code = "en" if output_language == "EN" else "fi"

    with results_col:
        st.markdown("""
            <div class="status-banner">
                Course Pack processed successfully. Each file was analyzed independently to protect quality and language direction.
            </div>
        """, unsafe_allow_html=True)

        st.markdown(f"""
            <div class="result-context">
                <span class="result-context-chip">Course: {html.escape(selected_course["name"])}</span>
                <span class="result-context-chip">Files uploaded: {len(batch_results)}</span>
                <span class="result-context-chip">Processed: {processed_count}</span>
                <span class="result-context-chip">Skipped / review: {skipped_count}</span>
            </div>
        """, unsafe_allow_html=True)

        st.markdown(f"""
            <div class="kpi-grid">
                <div class="kpi-tile accent-blue">
                    <p class="kpi-label">Files Uploaded</p>
                    <p class="kpi-value">{len(batch_results)}</p>
                    <p class="kpi-copy">Course files submitted to batch processing.</p>
                </div>
                <div class="kpi-tile accent-green">
                    <p class="kpi-label">Processed</p>
                    <p class="kpi-value">{processed_count}</p>
                    <p class="kpi-copy">Files that received a full Language Station output.</p>
                </div>
                <div class="kpi-tile">
                    <p class="kpi-label">English / Finnish</p>
                    <p class="kpi-value">{english_count} / {finnish_count}</p>
                    <p class="kpi-copy">Detected source-language split across the batch.</p>
                </div>
                <div class="kpi-tile accent-gold">
                    <p class="kpi-label">Needs Review</p>
                    <p class="kpi-value">{skipped_count}</p>
                    <p class="kpi-copy">Files skipped because extraction or language detection was not safe enough.</p>
                </div>
            </div>
        """, unsafe_allow_html=True)

        for item in batch_results:
            st.markdown(f"""
                <div class="result-context" style="margin-top: 1rem;">
                    <span class="result-context-chip">{html.escape(item["filename"])}</span>
                    <span class="result-context-chip">Type: {html.escape(item["file_type"].upper())}</span>
                    <span class="result-context-chip">Detected: {html.escape(item["source_language"])}</span>
                    <span class="result-context-chip">Direction: {html.escape(item["resolved_direction"] or "Skipped")}</span>
                    <span class="result-context-chip">Status: {html.escape(item["status"])}</span>
                </div>
            """, unsafe_allow_html=True)

            if item["warnings"]:
                for warning in item["warnings"]:
                    st.warning(warning)

            if item["status"] != "Processed" or not item["result"]:
                continue

            result = item["result"]
            source_text = item["source_text"]
            source_word_count = max(len(source_text.split()), 1)
            simplified_text_value = get_localized_text(result.simplified_text, output_language_code)
            simplified_word_count = len(simplified_text_value.split())

            if source_word_count < 100:
                simplification_delta = f"{simplified_word_count} w"
            else:
                expansion = round(((simplified_word_count / source_word_count) - 1) * 100)
                simplification_delta = f"+{expansion}%" if expansion >= 0 else f"-{abs(expansion)}%"

            tab_names = [f"{item['filename']} Terms", f"{item['filename']} Text"]
            if has_selected_students:
                tab_names.append(f"{item['filename']} Activities")

            tabs = st.tabs(tab_names)
            with tabs[0]:
                for glossary_item in result.glossary:
                    term = html.escape(glossary_item.term)
                    translation = html.escape(glossary_item.equivalent_term)
                    academic_definition = html.escape(get_localized_text(glossary_item.academic_definition, output_language_code))
                    simple_definition = html.escape(get_localized_text(glossary_item.simple_definition, output_language_code))
                    cognitive_note = html.escape(get_localized_text(glossary_item.cognitive_note, output_language_code))
                    trans_label = "Finnish equivalent:" if item["resolved_direction"] == "English → Finnish" else "English equivalent:"
                    if output_language_code == "fi":
                        trans_label = "Suomenkielinen vastine:" if item["resolved_direction"] == "English → Finnish" else "Englanninkielinen vastine:"

                    st.markdown(f"""
                    <div class="glossary-card">
                        <div class="card-meta">
                            <span class="term-chip">Verified Term</span>
                            <span class="translation-chip"><b>{trans_label}</b> <span style="font-size:16px;">{translation}</span></span>
                            <span class="difficulty-chip">High Cognitive Load</span>
                        </div>
                        <h3 class="term-title">{term}</h3>
                        <span class="card-label">Academic Definition</span>
                        <p class="card-copy">{academic_definition}</p>
                        <span class="card-label">Simplified Meaning</span>
                        <p class="card-copy">{simple_definition}</p>
                        <div class="cognitive-note">
                            <strong>Why it is difficult:</strong> {cognitive_note}
                        </div>
                    </div>
                    """, unsafe_allow_html=True)

            with tabs[1]:
                display_text = source_text[:15000] + "..." if len(source_text) > 15000 else source_text
                original_text = html.escape(display_text)
                simplified_text = html.escape(simplified_text_value)
                original_words = len(display_text.split())
                adapted_words = len(simplified_text_value.split())
                st.markdown(f"""
                    <div class="bridge-stats">
                        <div class="bridge-stat">
                            <p class="bridge-stat-label">Original Length</p>
                            <p class="bridge-stat-value">{original_words} words</p>
                        </div>
                        <div class="bridge-stat">
                            <p class="bridge-stat-label">Adapted Length</p>
                            <p class="bridge-stat-value">{adapted_words} words</p>
                        </div>
                        <div class="bridge-stat">
                            <p class="bridge-stat-label">Scaffolding Added</p>
                            <p class="bridge-stat-value">{simplification_delta}</p>
                        </div>
                    </div>
                    <div class="bridge-grid">
                        <div class="text-panel original">
                            <div class="panel-label original">Original Source</div>
                            <h3 class="panel-title">Dense academic wording</h3>
                            <p class="panel-copy">{original_text}</p>
                        </div>
                        <div class="text-panel simplified">
                            <div class="panel-label simplified">CEFR B1 Adaptation</div>
                            <h3 class="panel-title">Clearer language for multilingual learners</h3>
                            <p class="panel-copy">{simplified_text}</p>
                        </div>
                    </div>
                """, unsafe_allow_html=True)

            if has_selected_students:
                with tabs[2]:
                    for i, activity in enumerate(result.pedagogy_suggestions, 1):
                        activity_name = html.escape(get_localized_text(activity.activity_name, output_language_code))
                        activity_sections = format_activity_sections(
                            get_localized_text(activity.instructions, output_language_code),
                            output_language_code
                        )
                        goal = html.escape(activity_sections["goal"])
                        languages = html.escape(activity_sections["languages"])
                        why = html.escape(activity_sections["why"])
                        task = html.escape(activity_sections["task"])
                        step_items = "".join([
                            f"<li>{html.escape(step)}</li>" for step in activity_sections["step_list"]
                        ])
                        st.markdown(f"""
                        <div class="planner-card">
                            <div class="planner-chip">Activity {i}</div>
                            <h3 class="planner-title">{activity_name}</h3>
                            <div class="activity-grid">
                                <div class="activity-side">
                                    <div class="activity-block">
                                        <p class="activity-block-title">Goal</p>
                                        <p class="activity-block-copy">{goal}</p>
                                    </div>
                                    <div class="activity-block soft-green">
                                        <p class="activity-block-title">Languages Used</p>
                                        <p class="activity-block-copy">{languages}</p>
                                    </div>
                                    <div class="activity-block soft-gold">
                                        <p class="activity-block-title">Why This Works</p>
                                        <p class="activity-block-copy">{why}</p>
                                    </div>
                                </div>
                                <div class="activity-steps">
                                    <p class="activity-block-title">Implementation</p>
                                    <p class="activity-block-copy" style="margin-bottom: 0.75rem;">{task}</p>
                                    <ol class="activity-step-list">
                                        {step_items}
                                    </ol>
                                </div>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
elif st.session_state.result:
    result = st.session_state.result
    source_text = st.session_state.get('source_text_cache', "")
    source_processing_summary = st.session_state.get("source_processing_summary")
    has_selected_students = bool(selected_students)
    source_word_count = max(len(source_text.split()), 1)
    results_col, toggle_col = st.columns([0.92, 0.08], gap="medium")

    with toggle_col:
        st.markdown('<div class="floating-toggle-label">Output</div>', unsafe_allow_html=True)
        output_language = st.radio(
            "AI output language",
            ["EN", "FIN"],
            horizontal=False,
            label_visibility="collapsed",
            format_func=lambda x: {"EN": "🇬🇧", "FIN": "🇫🇮"}[x],
            key="output_ui_language"
        )

    output_language_code = "en" if output_language == "EN" else "fi"
    simplified_text_value = get_localized_text(result.simplified_text, output_language_code)
    simplified_word_count = len(simplified_text_value.split())
    
    if source_word_count < 100:
        simplification_delta = f"{simplified_word_count} w"
        simplification_label = "Material Generated"
        simplification_copy = "Length of the generated classroom material."
    else:
        expansion = round(((simplified_word_count / source_word_count) - 1) * 100)
        simplification_delta = f"+{expansion}%" if expansion >= 0 else f"-{abs(expansion)}%"
        simplification_label = "Scaffolding Added"
        simplification_copy = "Estimated text expansion required to safely scaffold dense academic concepts."
    
    with results_col:
        st.markdown("""
            <div class="status-banner">
                Language Station generated successfully. Your multilingual teaching kit is ready for review.
            </div>
        """, unsafe_allow_html=True)

        st.markdown(f"""
            <div class="result-context">
                <span class="result-context-chip">Course: {html.escape(selected_course["name"])}</span>
                <span class="result-context-chip">Group size: {len(selected_students)} students</span>
                <span class="result-context-chip">Direction: {html.escape(language_direction)}</span>
                <span class="result-context-chip">Input mode: {html.escape(input_mode)}</span>
            </div>
        """, unsafe_allow_html=True)

        st.info("This tool transforms monolingual materials into multilingual, translanguaging-based learning experiences.")

        if source_processing_summary and source_processing_summary.get("kind") == "pptx":
            notes_status = "Included" if source_processing_summary.get("has_notes") else "Not found"
            st.caption(
                "PowerPoint extraction summary: "
                f"{source_processing_summary.get('slides_processed', 0)} slides processed, "
                f"speaker notes {notes_status.lower()}."
            )
            if source_processing_summary.get("visual_warning"):
                st.warning("Some slides appear to rely on images, charts, or layout-heavy content that may not be fully captured from the PowerPoint.")

        st.markdown(f"""
            <div class="kpi-grid">
                <div class="kpi-tile accent-blue">
                    <p class="kpi-label">Academic Terms</p>
                    <p class="kpi-value">{len(result.glossary)}</p>
                    <p class="kpi-copy">Priority terminology identified for multilingual support.</p>
                </div>
                <div class="kpi-tile accent-green">
                    <p class="kpi-label">{simplification_label}</p>
                    <p class="kpi-value">{simplification_delta}</p>
                    <p class="kpi-copy">{simplification_copy}</p>
                </div>
                <div class="kpi-tile">
                    <p class="kpi-label">{'Learning Activities' if has_selected_students else 'Classroom Group'}</p>
                    <p class="kpi-value">{len(result.pedagogy_suggestions) if has_selected_students else len(selected_students)}</p>
                    <p class="kpi-copy">{'Ready-to-use classroom activities generated from the source.' if has_selected_students else 'Add students to generate collaborative translanguaging activities.'}</p>
                </div>
                <div class="kpi-tile accent-gold">
                    <p class="kpi-label">Delivery Time</p>
                    <p class="kpi-value">Seconds</p>
                    <p class="kpi-copy">Teacher preparation support delivered in a single generation cycle.</p>
                </div>
            </div>
        """, unsafe_allow_html=True)

        if has_selected_students:
            tab1, tab2, tab3 = st.tabs([
                "Key Academic Terms",
                "Accessible Text",
                "Learning Activities"
            ])
        else:
            tab1, tab2 = st.tabs([
                "Key Academic Terms",
                "Accessible Text"
            ])

        with tab1:
            st.markdown("""
                <div class="section-intro">
                    <div class="section-kicker">Academic Language Support</div>
                    <h2 class="section-title">Key Academic Terms</h2>
                    <p class="section-copy">
                        Verified academic terms with semantic scaffolding, and deterministic Finnish equivalents for high-stakes accuracy.
                    </p>
                </div>
            """, unsafe_allow_html=True)

            for item in result.glossary:
                term = html.escape(item.term)
                translation = html.escape(item.equivalent_term)
                academic_definition = html.escape(get_localized_text(item.academic_definition, output_language_code))
                simple_definition = html.escape(get_localized_text(item.simple_definition, output_language_code))
                cognitive_note = html.escape(get_localized_text(item.cognitive_note, output_language_code))

                # Dynamic Labeling based on Mode
                lang_dir = st.session_state.get('lang_dir_cache', "English → Finnish")
                if lang_dir == "English → Finnish":
                    trans_label = "Finnish equivalent:" if output_language == "EN" else "Suomenkielinen vastine:"
                else:
                    trans_label = "English equivalent:" if output_language == "EN" else "Englanninkielinen vastine:"

                st.markdown(f"""
                <div class="glossary-card">
                    <div class="card-meta">
                        <span class="term-chip">Verified Term</span>
                        <span class="translation-chip"><b>{trans_label}</b> <span style="font-size:16px;">{translation}</span></span>
                        <span class="difficulty-chip">High Cognitive Load</span>
                    </div>
                    <h3 class="term-title">{term}</h3>
                    <span class="card-label">Academic Definition</span>
                    <p class="card-copy">{academic_definition}</p>
                    <span class="card-label">Simplified Meaning</span>
                    <p class="card-copy">{simple_definition}</p>
                    <div class="cognitive-note">
                        <strong>Why it is difficult:</strong> {cognitive_note}
                    </div>
                </div>
                """, unsafe_allow_html=True)

        with tab2:
            st.markdown("""
                <div class="section-intro">
                    <div class="section-kicker">Before And After</div>
                    <h2 class="section-title">Bilingual Bridge</h2>
                    <p class="section-copy">
                        Compare the original academic source with a more accessible B1 adaptation while preserving the core teaching content.
                    </p>
                </div>
            """, unsafe_allow_html=True)
            if output_language == "EN":
                st.subheader("Accessible Version for Understanding and Discussion (CEFR B1)")
            else:
                st.subheader("Helposti luettava versio ymmartamista ja keskustelua varten (CEFR B1)")

            display_text = source_text[:15000] + "..." if len(source_text) > 15000 else source_text
            original_text = html.escape(display_text)
            simplified_text = html.escape(simplified_text_value)
            original_words = len(display_text.split())
            adapted_words = len(simplified_text_value.split())
            complexity_value = simplification_delta if source_word_count >= 100 else "Generated"
            st.markdown(f"""
                <div class="bridge-stats">
                    <div class="bridge-stat">
                        <p class="bridge-stat-label">Original Length</p>
                        <p class="bridge-stat-value">{original_words} words</p>
                    </div>
                    <div class="bridge-stat">
                        <p class="bridge-stat-label">Adapted Length</p>
                        <p class="bridge-stat-value">{adapted_words} words</p>
                    </div>
                    <div class="bridge-stat">
                        <p class="bridge-stat-label">Scaffolding Added</p>
                        <p class="bridge-stat-value">{complexity_value}</p>
                    </div>
                </div>
            """, unsafe_allow_html=True)
            st.markdown(f"""
                <div class="bridge-grid">
                    <div class="text-panel original">
                        <div class="panel-label original">Original Source</div>
                        <h3 class="panel-title">Dense academic wording</h3>
                        <p class="panel-copy">{original_text}</p>
                    </div>
                    <div class="text-panel simplified">
                        <div class="panel-label simplified">CEFR B1 Adaptation</div>
                        <h3 class="panel-title">Clearer language for multilingual learners</h3>
                        <p class="panel-copy">{simplified_text}</p>
                    </div>
                </div>
            """, unsafe_allow_html=True)

        if has_selected_students:
            with tab3:
                st.markdown(f"""
                    <div class="section-intro">
                        <div class="section-kicker">Teacher-Ready Outputs</div>
                        <h2 class="section-title">Pedagogy Planner</h2>
                        <p class="section-copy">
                            Text-specific classroom activities designed to support collaboration, translanguaging, and fast lesson preparation.
                        </p>
                    </div>
                """, unsafe_allow_html=True)
                st.markdown("## Lesson Planning: Collaborative Groupwork Activities")
                st.caption("Designed to support lesson planning with ready-to-use teaching methods.")
                st.markdown("### Ready-to-Use Classroom Activities")
                st.caption("Designed for immediate use in multilingual classrooms")

                for i, activity in enumerate(result.pedagogy_suggestions, 1):
                    activity_name = html.escape(get_localized_text(activity.activity_name, output_language_code))
                    activity_sections = format_activity_sections(
                        get_localized_text(activity.instructions, output_language_code),
                        output_language_code
                    )
                    goal = html.escape(activity_sections["goal"])
                    languages = html.escape(activity_sections["languages"])
                    why = html.escape(activity_sections["why"])
                    task = html.escape(activity_sections["task"])
                    step_items = "".join([
                        f"<li>{html.escape(step)}</li>" for step in activity_sections["step_list"]
                    ])
                    st.markdown(f"""
                    <div class="planner-card">
                        <div class="planner-chip">Activity {i}</div>
                        <h3 class="planner-title">{activity_name}</h3>
                        <div class="activity-grid">
                            <div class="activity-side">
                                <div class="activity-block">
                                    <p class="activity-block-title">Goal</p>
                                    <p class="activity-block-copy">{goal}</p>
                                </div>
                                <div class="activity-block soft-green">
                                    <p class="activity-block-title">Languages Used</p>
                                    <p class="activity-block-copy">{languages}</p>
                                </div>
                                <div class="activity-block soft-gold">
                                    <p class="activity-block-title">Why This Works</p>
                                    <p class="activity-block-copy">{why}</p>
                                </div>
                            </div>
                            <div class="activity-steps">
                                <p class="activity-block-title">Implementation</p>
                                <p class="activity-block-copy" style="margin-bottom: 0.75rem;">{task}</p>
                                <ol class="activity-step-list">
                                    {step_items}
                                </ol>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                st.caption("Copy and use these directly in your classroom.")
else:
    st.markdown("""
        <div class="empty-state">
            Choose your input mode in the sidebar to either upload academic material or describe a lesson plan.
            Then click <strong>Generate Language Station</strong> to instantly create a cognitive glossary,
            CEFR-adapted text, and ready-to-use pedagogical activities.
        </div>
    """, unsafe_allow_html=True)

# Footer
st.markdown("""<div class="footer-note">Developed for Sohjo Hacks | Multilingual Pedagogy Project</div>""", unsafe_allow_html=True)
